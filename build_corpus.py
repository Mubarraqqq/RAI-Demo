#!/usr/bin/env python3

import argparse
import json
import os
import re
import subprocess
import sys
from pathlib import Path
from typing import Any

try:
    from dotenv import load_dotenv
except ModuleNotFoundError:  # pragma: no cover - optional dependency
    def load_dotenv(*args, **kwargs):  # type: ignore[override]
        return False
from pptx import Presentation

from storage import Storage


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Build normalized documents from the corpus registry."
    )
    parser.add_argument(
        "--registry",
        default="data/processed/corpus_registry.json",
        help="Path to the canonical corpus registry JSON file.",
    )
    parser.add_argument(
        "--output",
        default="data/processed/documents.jsonl",
        help="Path to write normalized documents as JSON Lines.",
    )
    parser.add_argument(
        "--chunks-output",
        default="data/processed/chunks.jsonl",
        help="Path to write normalized chunks as JSON Lines.",
    )
    parser.add_argument(
        "--embeddings-output",
        default="data/processed/embeddings.jsonl",
        help="Path to write chunk embeddings as JSON Lines.",
    )
    parser.add_argument(
        "--embed",
        action="store_true",
        help="Generate OpenAI embeddings for chunk records.",
    )
    parser.add_argument(
        "--embedding-model",
        default="text-embedding-3-small",
        help="OpenAI embedding model to use when --embed is set.",
    )
    parser.add_argument(
        "--module-cache",
        default="data/processed/.swift-module-cache",
        help="Directory used by the Swift PDF fallback extractor.",
    )
    parser.add_argument(
        "--bootstrap-db",
        action="store_true",
        help="Load the generated documents, chunks, and optional embeddings into SQLite.",
    )
    parser.add_argument(
        "--db-path",
        default="data/processed/manifest.db",
        help="Path to the SQLite database used when --bootstrap-db is set.",
    )
    return parser.parse_args()


def load_registry(path: Path) -> dict[str, Any]:
    try:
        return json.loads(path.read_text())
    except FileNotFoundError as exc:
        raise SystemExit(f"Registry file not found: {path}") from exc
    except json.JSONDecodeError as exc:
        raise SystemExit(f"Registry file is not valid JSON: {path}") from exc


def normalize_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, max_chars: int = 900, overlap_paragraphs: int = 1) -> list[str]:
    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    if not paragraphs:
        return []

    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for paragraph in paragraphs:
        paragraph_len = len(paragraph)
        if current and current_len + paragraph_len + 2 > max_chars:
            chunks.append("\n\n".join(current))
            overlap = current[-overlap_paragraphs:] if overlap_paragraphs else []
            current = overlap.copy()
            current_len = sum(len(item) for item in current) + max(len(current) - 1, 0) * 2

        current.append(paragraph)
        current_len += paragraph_len + (2 if len(current) > 1 else 0)

    if current:
        chunks.append("\n\n".join(current))

    return [normalize_text(chunk) for chunk in chunks if chunk.strip()]


def extract_pdf_text(path: Path, module_cache: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        parts = [page.extract_text() or "" for page in reader.pages]
        return normalize_text("\n\n".join(parts))
    except ModuleNotFoundError:
        return extract_pdf_text_with_swift(path, module_cache)


def extract_pdf_text_with_swift(path: Path, module_cache: Path) -> str:
    module_cache.mkdir(parents=True, exist_ok=True)
    swift_script = f"""
import Foundation
import PDFKit

let url = URL(fileURLWithPath: {json.dumps(str(path.resolve()))})
guard let doc = PDFDocument(url: url) else {{
    fputs("Failed to open PDF: {path}\\n", stderr)
    exit(1)
}}

for index in 0..<doc.pageCount {{
    if let pageText = doc.page(at: index)?.string {{
        print(pageText)
        print("\\n__PAGE_BREAK__\\n")
    }}
}}
"""
    cmd = [
        "swift",
        "-module-cache-path",
        str(module_cache),
        "-",
    ]
    result = subprocess.run(
        cmd,
        input=swift_script,
        text=True,
        capture_output=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(
            f"PDF extraction failed for {path}: {result.stderr.strip() or result.stdout.strip()}"
        )
    text = result.stdout.replace("\n__PAGE_BREAK__\n", "\n\n")
    return normalize_text(text)


def extract_docx_text(path: Path) -> str:
    try:
        import docx  # type: ignore

        document = docx.Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs]
        return normalize_text("\n".join(parts))
    except ModuleNotFoundError:
        return extract_docx_text_with_textutil(path)


def extract_docx_text_with_textutil(path: Path) -> str:
    cmd = ["textutil", "-convert", "txt", "-stdout", str(path)]
    result = subprocess.run(
        cmd,
        text=True,
        capture_output=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(
            f"DOCX extraction failed for {path}: {result.stderr.strip() or result.stdout.strip()}"
        )
    return normalize_text(result.stdout)


def extract_pptx_text(path: Path) -> str:
    presentation = Presentation(str(path))
    slides_text: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if parts:
            slides_text.append(f"Slide {index}\n" + "\n".join(parts))
    return normalize_text("\n\n".join(slides_text))


def extract_text(path: Path, module_cache: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path, module_cache)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    raise RuntimeError(f"Unsupported file type: {path.suffix or '<no extension>'}")


def build_document_record(
    subfacet: dict[str, Any],
    asset_type: str,
    asset_index: int,
    asset_path: Path,
    text: str,
) -> dict[str, Any]:
    document_id = f"{subfacet['subfacet_id']}::{asset_type}::{asset_index}"
    return {
        "document_id": document_id,
        "subfacet_id": subfacet["subfacet_id"],
        "canonical_name": subfacet["canonical_name"],
        "facet": subfacet["facet"],
        "asset_type": asset_type,
        "source_path": str(asset_path),
        "source_filename": asset_path.name,
        "extension": asset_path.suffix.lower(),
        "text": text,
        "text_length": len(text),
    }


def build_chunk_records(documents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    chunk_records: list[dict[str, Any]] = []
    for document in documents:
        chunks = chunk_text(document["text"])
        for index, chunk in enumerate(chunks, start=1):
            chunk_records.append(
                {
                    "chunk_id": f"{document['document_id']}::chunk::{index}",
                    "document_id": document["document_id"],
                    "subfacet_id": document["subfacet_id"],
                    "canonical_name": document["canonical_name"],
                    "facet": document["facet"],
                    "asset_type": document["asset_type"],
                    "source_path": document["source_path"],
                    "source_filename": document["source_filename"],
                    "chunk_index": index,
                    "text": chunk,
                    "text_length": len(chunk),
                }
            )
    return chunk_records


def write_jsonl(path: Path, records: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as handle:
        for record in records:
            handle.write(json.dumps(record, ensure_ascii=True) + "\n")


def build_embedding_records(
    chunks: list[dict[str, Any]],
    model: str = "text-embedding-3-small",
    batch_size: int = 64,
) -> list[dict[str, Any]]:
    if not os.environ.get("OPENAI_API_KEY"):
        raise RuntimeError("OPENAI_API_KEY must be set to build OpenAI embeddings.")

    try:
        from openai import OpenAI
    except ModuleNotFoundError as exc:
        raise RuntimeError("The openai package is required to build embeddings.") from exc

    client = OpenAI()
    records: list[dict[str, Any]] = []
    for start in range(0, len(chunks), batch_size):
        batch = chunks[start : start + batch_size]
        response = client.embeddings.create(
            model=model,
            input=[chunk["text"] for chunk in batch],
        )
        for chunk, embedding_data in zip(batch, response.data):
            records.append(
                {
                    "chunk_id": chunk["chunk_id"],
                    "model": model,
                    "embedding": embedding_data.embedding,
                }
            )
    return records


def iter_assets(subfacet: dict[str, Any]) -> list[tuple[str, list[str]]]:
    return [
        ("slides", subfacet.get("slide_assets", [])),
        ("worksheets", subfacet.get("worksheet_assets", [])),
        ("transcripts", subfacet.get("transcript_assets", [])),
        ("posts", subfacet.get("post_assets", [])),
    ]


def build_documents(
    registry: dict[str, Any],
    output_path: Path,
    module_cache: Path,
) -> tuple[int, int]:
    documents: list[dict[str, Any]] = []
    failures = 0

    for subfacet in registry.get("subfacets", []):
        for asset_type, asset_paths in iter_assets(subfacet):
            for asset_index, asset_path_str in enumerate(asset_paths, start=1):
                asset_path = Path(asset_path_str)
                if not asset_path.exists():
                    failures += 1
                    print(
                        f"[warn] Missing {asset_type} file for {subfacet['canonical_name']}: {asset_path}",
                        file=sys.stderr,
                    )
                    continue
                try:
                    text = extract_text(asset_path, module_cache)
                except Exception as exc:  # noqa: BLE001
                    failures += 1
                    print(
                        f"[warn] Failed to extract {asset_type} file for {subfacet['canonical_name']}: "
                        f"{asset_path} ({exc})",
                        file=sys.stderr,
                    )
                    continue

                documents.append(
                    build_document_record(
                        subfacet=subfacet,
                        asset_type=asset_type,
                        asset_index=asset_index,
                        asset_path=asset_path,
                        text=text,
                    )
                )

    write_jsonl(output_path, documents)

    return len(documents), failures


def main() -> int:
    load_dotenv()

    args = parse_args()
    registry_path = Path(args.registry)
    output_path = Path(args.output)
    chunks_output_path = Path(args.chunks_output)
    embeddings_output_path = Path(args.embeddings_output)
    module_cache = Path(args.module_cache)

    registry = load_registry(registry_path)
    document_count, failure_count = build_documents(
        registry=registry,
        output_path=output_path,
        module_cache=module_cache,
    )
    documents = [
        json.loads(line)
        for line in output_path.read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]
    chunk_records = build_chunk_records(documents)
    write_jsonl(chunks_output_path, chunk_records)
    print(f"Wrote {document_count} documents to {output_path}")
    print(f"Wrote {len(chunk_records)} chunks to {chunks_output_path}")
    if args.embed:
        embedding_records = build_embedding_records(
            chunk_records,
            model=args.embedding_model,
        )
        write_jsonl(embeddings_output_path, embedding_records)
        print(
            f"Wrote {len(embedding_records)} embeddings to {embeddings_output_path} "
            f"using {args.embedding_model}"
        )

    if args.bootstrap_db:
        store = Storage(args.db_path)
        store.initialize()
        loaded_documents = store.load_documents(str(output_path))
        loaded_chunks = store.load_chunks(str(chunks_output_path))
        print(f"Loaded {loaded_documents} documents into {args.db_path}")
        print(f"Loaded {loaded_chunks} chunks into {args.db_path}")
        if args.embed:
            loaded_embeddings = store.load_embeddings(str(embeddings_output_path))
            print(f"Loaded {loaded_embeddings} embeddings into {args.db_path}")

    if failure_count:
        print(f"Encountered {failure_count} extraction warning(s)", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
